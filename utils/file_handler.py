from pathlib import Path
from tempfile import NamedTemporaryFile

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from werkzeug.utils import secure_filename


ALLOWED_DOCUMENT_EXTENSIONS = {".pdf", ".docx", ".txt", ".pptx"}
MAX_FILE_SIZE_MB = 100
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024


def trim_to_word_boundary(text, max_characters):
    """
    Trim extracted text without cutting the final word in half.
    """
    if not max_characters or len(text) <= max_characters:
        return text

    trimmed_text = text[:max_characters].rsplit(" ", 1)[0].strip()
    return trimmed_text or text[:max_characters].strip()


def append_limited_section(sections, text, current_length, max_characters=None):
    """
    Append extracted text and stop once the caller's character limit is reached.
    """
    section = (text or "").strip()
    if not section:
        return current_length, False

    separator_length = 2 if sections else 0
    if max_characters:
        remaining = max_characters - current_length - separator_length
        if remaining <= 0:
            return current_length, True

        if len(section) > remaining:
            section = trim_to_word_boundary(section, remaining)
            if section:
                sections.append(section)
                current_length += separator_length + len(section)
            return current_length, True

    sections.append(section)
    current_length += separator_length + len(section)
    return current_length, bool(max_characters and current_length >= max_characters)


def validate_document_extension(filename):
    """
    Ensure the uploaded document uses a supported extension.
    """
    suffix = Path(filename or "").suffix.lower()
    if suffix not in ALLOWED_DOCUMENT_EXTENSIONS:
        raise ValueError("Only PDF, DOCX, TXT, and PPTX files are supported")
    return suffix


def save_uploaded_document_temporarily(file_storage, temp_dir):
    """
    Save an uploaded document into the temporary uploads directory.
    """
    if file_storage is None:
        raise ValueError("No file was provided")

    filename = (file_storage.filename or "").strip()
    if not filename:
        raise ValueError("Please choose a file")

    suffix = validate_document_extension(filename)
    temp_dir = Path(temp_dir)
    temp_dir.mkdir(parents=True, exist_ok=True)

    safe_stem = secure_filename(Path(filename).stem) or "upload"
    temp_file = NamedTemporaryFile(
        delete=False,
        dir=temp_dir,
        prefix=f"{safe_stem}_",
        suffix=suffix,
    )
    temp_path = Path(temp_file.name)
    temp_file.close()

    file_storage.save(temp_path)

    if temp_path.stat().st_size == 0:
        temp_path.unlink(missing_ok=True)
        raise ValueError("Uploaded file is empty")

    # Check file size limit (100MB)
    if temp_path.stat().st_size > MAX_FILE_SIZE_BYTES:
        temp_path.unlink(missing_ok=True)
        raise ValueError(f"File size exceeds {MAX_FILE_SIZE_MB}MB limit")

    return temp_path


def clean_extracted_text(text):
    """
    Trim noisy whitespace while preserving paragraph breaks.
    """
    compact_lines = []
    last_line_blank = False

    for raw_line in text.replace("\x00", "").splitlines():
        line = raw_line.strip()
        if line:
            compact_lines.append(line)
            last_line_blank = False
        elif not last_line_blank:
            compact_lines.append("")
            last_line_blank = True

    cleaned_text = "\n".join(compact_lines).strip()
    if not cleaned_text:
        raise ValueError("No readable text found in the uploaded file")

    return cleaned_text


def extract_text_from_pdf(file_path, max_characters=None):
    """
    Extract readable text from a PDF document.
    """
    try:
        reader = PdfReader(str(file_path))
    except Exception as exc:
        raise ValueError("Unable to read the uploaded PDF file") from exc

    page_text = []
    current_length = 0
    for page in reader.pages:
        extracted_text = page.extract_text() or ""
        current_length, is_limit_reached = append_limited_section(
            page_text,
            extracted_text,
            current_length,
            max_characters=max_characters,
        )
        if is_limit_reached:
            break

    return clean_extracted_text("\n\n".join(page_text))


def extract_text_from_docx(file_path, max_characters=None):
    """
    Extract text from DOCX paragraphs and tables.
    """
    try:
        document = Document(str(file_path))
    except Exception as exc:
        raise ValueError("Unable to read the uploaded DOCX file") from exc

    sections = []
    current_length = 0

    for paragraph in document.paragraphs:
        current_length, is_limit_reached = append_limited_section(
            sections,
            paragraph.text,
            current_length,
            max_characters=max_characters,
        )
        if is_limit_reached:
            return clean_extracted_text("\n\n".join(sections))

    for table in document.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if row_text:
                current_length, is_limit_reached = append_limited_section(
                    sections,
                    " | ".join(row_text),
                    current_length,
                    max_characters=max_characters,
                )
                if is_limit_reached:
                    return clean_extracted_text("\n\n".join(sections))

    return clean_extracted_text("\n\n".join(sections))


def extract_text_from_txt(file_path, max_characters=None):
    """
    Read plain text from a TXT file.
    """
    path = Path(file_path)
    read_size = max_characters + 4096 if max_characters else None

    try:
        with path.open("r", encoding="utf-8-sig") as file:
            text = file.read(read_size)
    except UnicodeDecodeError:
        try:
            with path.open("r", encoding="latin-1") as file:
                text = file.read(read_size)
        except Exception as exc:
            raise ValueError("Unable to read the uploaded TXT file") from exc
    except Exception as exc:
        raise ValueError("Unable to read the uploaded TXT file") from exc

    return clean_extracted_text(trim_to_word_boundary(text, max_characters))


def extract_text_from_pptx(file_path, max_characters=None):
    """
    Extract text from PowerPoint presentations (slides and speaker notes).
    """
    try:
        presentation = Presentation(str(file_path))
    except Exception as exc:
        raise ValueError("Unable to read the uploaded PPTX file") from exc

    sections = []
    current_length = 0

    # Extract text from all slides and shapes
    for slide in presentation.slides:
        for shape in slide.shapes:
            # Extract text from text frames
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                current_length, is_limit_reached = append_limited_section(
                    sections,
                    shape.text,
                    current_length,
                    max_characters=max_characters,
                )
                if is_limit_reached:
                    return clean_extracted_text("\n\n".join(sections))

            # Extract text from tables within shapes
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                    if row_text:
                        current_length, is_limit_reached = append_limited_section(
                            sections,
                            " | ".join(row_text),
                            current_length,
                            max_characters=max_characters,
                        )
                        if is_limit_reached:
                            return clean_extracted_text("\n\n".join(sections))

        # Extract speaker notes if available
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                current_length, is_limit_reached = append_limited_section(
                    sections,
                    f"[Notes: {notes_text}]",
                    current_length,
                    max_characters=max_characters,
                )
                if is_limit_reached:
                    return clean_extracted_text("\n\n".join(sections))

    return clean_extracted_text("\n\n".join(sections))


def extract_text_from_file(file_path, max_characters=None):
    """
    Dispatch text extraction based on the uploaded document extension.
    """
    path = Path(file_path)
    suffix = path.suffix.lower()

    if suffix == ".pdf":
        return extract_text_from_pdf(path, max_characters=max_characters)
    if suffix == ".docx":
        return extract_text_from_docx(path, max_characters=max_characters)
    if suffix == ".txt":
        return extract_text_from_txt(path, max_characters=max_characters)
    if suffix == ".pptx":
        return extract_text_from_pptx(path, max_characters=max_characters)

    raise ValueError("Only PDF, DOCX, TXT, and PPTX files are supported")
